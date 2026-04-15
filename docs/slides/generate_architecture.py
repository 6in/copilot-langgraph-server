"""Copilot LangGraph Chat アーキテクチャ説明資料を .pptx として生成する。

実行:
    uv run --with python-pptx python docs/slides/generate_architecture.py

出力:
    docs/slides/architecture.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUTPUT = Path(__file__).parent / "architecture.pptx"

# ---- テーマカラー ----
C_BG = RGBColor(0xF7, 0xF9, 0xFC)
C_PRIMARY = RGBColor(0x1F, 0x4E, 0x79)
C_ACCENT = RGBColor(0x2E, 0x86, 0xAB)
C_OK = RGBColor(0x4C, 0xAF, 0x50)
C_WARN = RGBColor(0xE6, 0x8A, 0x00)
C_TEXT = RGBColor(0x22, 0x2B, 0x3A)
C_MUTED = RGBColor(0x6B, 0x73, 0x80)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BORDER = RGBColor(0xCF, 0xD8, 0xE3)


def _set_fill(shape, rgb: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _set_line(shape, rgb: RGBColor, width: float = 1.0) -> None:
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width)


def _set_text(
    tf,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = C_TEXT,
    align: int = PP_ALIGN.LEFT,
) -> None:
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.runs[0] if len(p.runs) > 0 else p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Meiryo"


def _add_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    fill: RGBColor = C_WHITE,
    border: RGBColor = C_BORDER,
    font_color: RGBColor = C_TEXT,
    size: int = 12,
    bold: bool = False,
    shape: int = MSO_SHAPE.ROUNDED_RECTANGLE,
):
    box = slide.shapes.add_shape(shape, Inches(left), Inches(top), Inches(width), Inches(height))
    _set_fill(box, fill)
    _set_line(box, border, 1.25)
    tf = box.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    _set_text(tf, text, size=size, bold=bold, color=font_color, align=PP_ALIGN.CENTER)
    return box


def _add_text(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = C_TEXT,
    align: int = PP_ALIGN.LEFT,
):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    _set_text(tf, text, size=size, bold=bold, color=color, align=align)
    return tb


def _add_bullets(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    items: list[str],
    *,
    size: int = 16,
    color: RGBColor = C_TEXT,
):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Meiryo"
    return tb


def _add_arrow(slide, x1: float, y1: float, x2: float, y2: float, *, color: RGBColor = C_ACCENT) -> None:
    conn = slide.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(2)
    # 矢印ヘッド
    line = conn.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    from lxml import etree

    tail = etree.SubElement(line, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")


def _add_section_header(slide, title: str, subtitle: str = "") -> None:
    # 左ストライプ
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.35), Inches(7.5))
    _set_fill(bar, C_PRIMARY)
    bar.line.fill.background()
    _add_text(slide, 0.6, 0.35, 12, 0.7, title, size=28, bold=True, color=C_PRIMARY)
    if subtitle:
        _add_text(slide, 0.6, 1.0, 12, 0.5, subtitle, size=14, color=C_MUTED)


def _blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    _set_fill(bg, C_BG)
    bg.line.fill.background()
    # 最背面に
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return slide


# -----------------------------------------------------------------------------
# スライド定義
# -----------------------------------------------------------------------------


def slide_title(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    # タイトル帯
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), Inches(13.33), Inches(2.5))
    _set_fill(band, C_PRIMARY)
    band.line.fill.background()
    _add_text(
        slide,
        0.8,
        2.9,
        11.7,
        1.2,
        "Copilot LangGraph Chat",
        size=44,
        bold=True,
        color=C_WHITE,
    )
    _add_text(
        slide,
        0.8,
        3.9,
        11.7,
        0.6,
        "GitHub Copilot を LangGraph で束ねる社内向け汎用チャット基盤",
        size=20,
        color=C_WHITE,
    )
    _add_text(slide, 0.8, 6.5, 11.7, 0.4, "アーキテクチャ説明資料  /  2026-04", size=14, color=C_MUTED)


def slide_overview(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_section_header(slide, "プロジェクト概要", "何を作っているか / なぜ Copilot SDK なのか")
    items = [
        "GitHub Copilot の JSON-RPC SDK を LangChain 互換の BaseChatModel としてラップ",
        "LangGraph で会話グラフ・ReAct ループ・Checkpointer 永続化を実装",
        "社内利用 200 名規模を想定したマルチユーザー・マルチアプリ構成",
        "5 つのアプリを同居: Chat / SuperChat / Gems / Canvas / DebateChat",
        "認証は GitHub Device Flow のみ。JWT HS256 を httpOnly cookie で管理",
    ]
    _add_bullets(slide, 0.8, 1.8, 12, 5, items, size=18)


def slide_apps(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_section_header(slide, "5 つのアプリケーション", "同じバックエンド基盤で切り替え可能")
    apps = [
        ("Chat", "シンプル ReAct チャット", C_ACCENT),
        ("SuperChat", "Router → SubAgent の\nオーケストレーション", C_PRIMARY),
        ("Gems", "システムプロンプト駆動の\nカスタム AI", C_OK),
        ("Canvas", "iframe RPC でホスト済み\nミニアプリと対話", C_WARN),
        ("DebateChat", "複数エージェントの\n議論シミュレーション", RGBColor(0x8E, 0x44, 0xAD)),
    ]
    x = 0.6
    for name, desc, color in apps:
        _add_box(slide, x, 2.3, 2.4, 2.4, name, fill=color, border=color, font_color=C_WHITE, size=20, bold=True)
        _add_box(slide, x, 4.8, 2.4, 1.8, desc, fill=C_WHITE, border=C_BORDER, size=12)
        x += 2.55


def slide_architecture(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_section_header(slide, "全体アーキテクチャ", "docker compose で起動するサービス構成")

    # Frontend
    _add_box(slide, 0.6, 2.0, 2.8, 1.0, "Frontend\nReact 19 + Vite (Bun)", fill=C_ACCENT, border=C_ACCENT, font_color=C_WHITE, size=14, bold=True)
    # API
    _add_box(slide, 4.4, 2.0, 2.8, 1.0, "FastAPI\n/api + SSE", fill=C_PRIMARY, border=C_PRIMARY, font_color=C_WHITE, size=14, bold=True)
    # Worker
    _add_box(slide, 8.2, 2.0, 2.8, 1.0, "arq Worker\nLangGraph 実行", fill=C_PRIMARY, border=C_PRIMARY, font_color=C_WHITE, size=14, bold=True)
    # MCP
    _add_box(slide, 11.0, 2.0, 2.1, 1.0, "MCP Server\nFastMCP", fill=C_WARN, border=C_WARN, font_color=C_WHITE, size=13, bold=True)

    # Redis
    _add_box(slide, 4.4, 3.8, 2.8, 0.9, "Redis  (arq queue / SSE notifier)", fill=C_WHITE, border=C_BORDER, size=13)
    # Postgres
    _add_box(slide, 8.2, 3.8, 2.8, 0.9, "PostgreSQL  (checkpointer)", fill=C_WHITE, border=C_BORDER, size=13)

    # Copilot SDK
    _add_box(slide, 4.4, 5.2, 6.6, 0.9, "github-copilot-sdk  (Device Flow / JSON-RPC)", fill=C_OK, border=C_OK, font_color=C_WHITE, size=14, bold=True)

    # 矢印
    _add_arrow(slide, 3.4, 2.5, 4.4, 2.5)  # FE -> API
    _add_arrow(slide, 7.2, 2.5, 8.2, 2.5)  # API -> Worker
    _add_arrow(slide, 11.0, 2.5, 11.05, 2.5)  # Worker -> MCP (隣接)
    _add_arrow(slide, 5.8, 3.0, 5.8, 3.8, color=C_MUTED)  # API -> Redis
    _add_arrow(slide, 9.6, 3.0, 9.6, 3.8, color=C_MUTED)  # Worker -> Postgres
    _add_arrow(slide, 9.6, 3.0, 5.8, 3.8, color=C_MUTED)  # Worker -> Redis
    _add_arrow(slide, 9.6, 4.7, 7.7, 5.2, color=C_MUTED)  # Worker -> Copilot
    _add_arrow(slide, 5.8, 3.0, 7.7, 5.2, color=C_MUTED)  # API -> Copilot (auth)

    _add_text(
        slide,
        0.8,
        6.5,
        12,
        0.8,
        "※ nginx で APP_PREFIX (/orochi) を strip して FastAPI に転送。Vite dev server (5173) が /api を API にプロキシ。",
        size=12,
        color=C_MUTED,
    )


def slide_auth_flow(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_section_header(slide, "認証フロー", "GitHub Device Flow + JWT (HS256) httpOnly cookie")

    steps = [
        ("1", "Browser", "POST /api/auth/device/start"),
        ("2", "FastAPI", "CopilotAuthManager\nが device code 発行"),
        ("3", "User", "GitHub で\ncode を承認"),
        ("4", "FastAPI", "poll → access_token\n→ 暗号化して保存"),
        ("5", "FastAPI", "JWT 発行\n→ httpOnly cookie"),
    ]
    x = 0.6
    for num, role, desc in steps:
        _add_box(slide, x, 2.2, 2.4, 0.7, f"Step {num} — {role}", fill=C_PRIMARY, border=C_PRIMARY, font_color=C_WHITE, size=13, bold=True)
        _add_box(slide, x, 2.95, 2.4, 1.6, desc, fill=C_WHITE, border=C_BORDER, size=12)
        if num != "5":
            _add_arrow(slide, x + 2.4, 3.75, x + 2.55, 3.75)
        x += 2.55

    _add_bullets(
        slide,
        0.8,
        5.0,
        12,
        2.0,
        [
            "Copilot access_token は Fernet で暗号化して DB 永続化（再ログイン不要）",
            "JWT は HS256、JTI ブロックリストで失効管理",
            "PAT 方式は対象外 — 非インタラクティブは今後の課題",
        ],
        size=15,
    )


def slide_async_job(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_section_header(slide, "非同期ジョブ処理", "POST → job_id → Worker 実行 → SSE 通知")

    # タイムライン
    _add_box(slide, 0.6, 2.2, 3.0, 1.0, "① Client\nPOST /api/chat", fill=C_ACCENT, border=C_ACCENT, font_color=C_WHITE, size=13, bold=True)
    _add_box(slide, 3.9, 2.2, 3.0, 1.0, "② FastAPI\narq.enqueue → job_id", fill=C_PRIMARY, border=C_PRIMARY, font_color=C_WHITE, size=13, bold=True)
    _add_box(slide, 7.2, 2.2, 3.0, 1.0, "③ Redis queue", fill=C_WHITE, border=C_BORDER, size=13)
    _add_box(slide, 10.5, 2.2, 2.6, 1.0, "④ arq Worker\nprocess_chat", fill=C_PRIMARY, border=C_PRIMARY, font_color=C_WHITE, size=13, bold=True)

    _add_arrow(slide, 3.6, 2.7, 3.9, 2.7)
    _add_arrow(slide, 6.9, 2.7, 7.2, 2.7)
    _add_arrow(slide, 10.2, 2.7, 10.5, 2.7)

    _add_box(slide, 10.5, 3.6, 2.6, 0.9, "⑤ JobStore に\n結果を格納", fill=C_OK, border=C_OK, font_color=C_WHITE, size=12, bold=True)
    _add_arrow(slide, 11.8, 3.2, 11.8, 3.6)

    _add_box(slide, 0.6, 4.8, 6.3, 1.0, "⑥ Client: GET /api/job/{id}/stream (SSE)", fill=C_WHITE, border=C_BORDER, size=14)
    _add_box(slide, 7.2, 4.8, 5.9, 1.0, "⑦ Notifier が完了イベント push", fill=C_WARN, border=C_WARN, font_color=C_WHITE, size=14, bold=True)
    _add_arrow(slide, 6.9, 5.3, 7.2, 5.3)

    _add_bullets(
        slide,
        0.8,
        6.0,
        12,
        1.2,
        [
            "WebSocket は使わず SSE のみ — 社内 nginx との相性と実装単純化を優先",
            "JobStore は In-memory + asyncio.Queue — 200 名規模では十分",
        ],
        size=14,
    )


def slide_langgraph(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_section_header(slide, "LangGraph + ChatCopilot", "ChatCopilot を BaseChatModel として差し込む")

    # グラフ図
    _add_box(slide, 1.0, 2.2, 2.6, 1.1, "START", fill=C_MUTED, border=C_MUTED, font_color=C_WHITE, size=14, bold=True, shape=MSO_SHAPE.OVAL)
    _add_box(slide, 4.3, 2.2, 2.6, 1.1, "agent\n(ChatCopilot)", fill=C_PRIMARY, border=C_PRIMARY, font_color=C_WHITE, size=13, bold=True)
    _add_box(slide, 7.6, 2.2, 2.6, 1.1, "ToolNode\n(LangChain tools)", fill=C_ACCENT, border=C_ACCENT, font_color=C_WHITE, size=13, bold=True)
    _add_box(slide, 10.9, 2.2, 2.2, 1.1, "END", fill=C_MUTED, border=C_MUTED, font_color=C_WHITE, size=14, bold=True, shape=MSO_SHAPE.OVAL)

    _add_arrow(slide, 3.6, 2.75, 4.3, 2.75)
    _add_arrow(slide, 6.9, 2.75, 7.6, 2.75)
    _add_arrow(slide, 10.2, 2.75, 10.9, 2.75)
    # ループバック
    _add_arrow(slide, 8.9, 3.3, 5.6, 3.3, color=C_MUTED)
    _add_text(slide, 6.8, 3.35, 2.0, 0.4, "tool_calls 残", size=11, color=C_MUTED)

    _add_bullets(
        slide,
        0.8,
        4.8,
        12,
        2.8,
        [
            "ChatCopilot: Copilot JSON-RPC を BaseChatModel にラップ (providers/copilot.py)",
            "BoundChatCopilot: bind_tools() でツールスキーマを system prompt に注入 → JSON 応答を AIMessage(tool_calls=[...]) に変換",
            "Checkpointer: PostgreSQL AsyncConnectionPool でスレッド永続化",
            "グラフは lifespan でコンパイル、ライフサイクルは caller が管理",
        ],
        size=15,
    )


def slide_subagent(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_section_header(slide, "SubAgent と MCP 統合", "AGENT.md から動的にエージェント生成")

    _add_box(slide, 0.6, 2.0, 12, 1.1, "SubAgentRegistry", fill=C_PRIMARY, border=C_PRIMARY, font_color=C_WHITE, size=18, bold=True)
    _add_text(slide, 0.8, 3.15, 12, 0.5, "agents/*/AGENT.md を自動ロード → tools: flag + mcp_tools で ToolEnabledSubAgent を生成", size=13, color=C_MUTED)

    _add_box(slide, 0.6, 4.0, 4.0, 1.8, "Router\n(SuperChat)\n↓ 宛先決定", fill=C_ACCENT, border=C_ACCENT, font_color=C_WHITE, size=14, bold=True)
    _add_box(slide, 5.0, 4.0, 4.0, 1.8, "SubAgent\n(ReAct mini graph)\nagent → ToolNode → agent", fill=C_PRIMARY, border=C_PRIMARY, font_color=C_WHITE, size=13, bold=True)
    _add_box(slide, 9.4, 4.0, 3.7, 1.8, "MCP Singleton\n(worker.startup)\nMultiServerMCPClient", fill=C_WARN, border=C_WARN, font_color=C_WHITE, size=13, bold=True)

    _add_arrow(slide, 4.6, 4.9, 5.0, 4.9)
    _add_arrow(slide, 9.0, 4.9, 9.4, 4.9)

    _add_bullets(
        slide,
        0.8,
        6.1,
        12,
        1.2,
        [
            "langchain-mcp-adapters で MCP tools → LangChain BaseTool に変換",
            "MCP 接続失敗時は ctx['mcp_tools']=[] で DEGRADED 継続 — 起動を止めない",
        ],
        size=14,
    )


def slide_persistence(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_section_header(slide, "永続化レイヤー", "状態はどこに置かれているか")

    _add_box(slide, 0.6, 2.2, 6.0, 2.5, "PostgreSQL", fill=C_PRIMARY, border=C_PRIMARY, font_color=C_WHITE, size=20, bold=True)
    _add_box(slide, 0.9, 3.2, 5.4, 0.6, "LangGraph Checkpointer (threads)", fill=C_WHITE, border=C_BORDER, size=12)
    _add_box(slide, 0.9, 3.85, 5.4, 0.6, "ユーザー / アプリ / Gem メタデータ", fill=C_WHITE, border=C_BORDER, size=12)
    _add_box(slide, 0.9, 4.0, 5.4, 0.6, "JWT JTI blocklist / 暗号化 token", fill=C_WHITE, border=C_BORDER, size=12)

    _add_box(slide, 7.0, 2.2, 6.0, 2.5, "Redis", fill=C_ACCENT, border=C_ACCENT, font_color=C_WHITE, size=20, bold=True)
    _add_box(slide, 7.3, 3.2, 5.4, 0.6, "arq ジョブキュー", fill=C_WHITE, border=C_BORDER, size=12)
    _add_box(slide, 7.3, 3.85, 5.4, 0.6, "SSE 通知用 asyncio.Queue ブリッジ", fill=C_WHITE, border=C_BORDER, size=12)

    _add_bullets(
        slide,
        0.8,
        5.2,
        12,
        2.0,
        [
            "Checkpointer は AsyncConnectionPool で共有 — lifespan で初期化",
            "config/db_pools.yaml で接続プールをチューニング (Phase 23)",
            "shared volume claude-code-outputs で worker / mcp-server 間の大容量出力を受け渡し",
        ],
        size=14,
    )


def slide_future(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_section_header(slide, "今後の拡張ポイント", "ロードマップ上の主なテーマ")
    _add_bullets(
        slide,
        0.8,
        2.0,
        12,
        5,
        [
            "data-ai-role 属性の導入 — AI エージェントが操作しやすい画面構成 (Phase 候補)",
            "Gem データの Redis キャッシュ (OrchestratorHandler)",
            "Slack Bot 連携 — 非インタラクティブ環境向けのエントリポイント追加検討",
            "Canvas: iframe RPC のモデル指定拡張と多段 AI 呼び出し",
            "監査ログ強化 — 200 名規模の利用履歴を追跡可能に",
        ],
        size=17,
    )


def slide_closing(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.0), Inches(13.33), Inches(1.5))
    _set_fill(band, C_PRIMARY)
    band.line.fill.background()
    _add_text(slide, 0.8, 3.3, 11.7, 1.0, "Questions & Discussion", size=40, bold=True, color=C_WHITE)
    _add_text(slide, 0.8, 5.0, 11.7, 0.6, "Source: docs/slides/generate_architecture.py", size=14, color=C_MUTED)


# -----------------------------------------------------------------------------
# Main
# -----------------------------------------------------------------------------


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_overview(prs)
    slide_apps(prs)
    slide_architecture(prs)
    slide_auth_flow(prs)
    slide_async_job(prs)
    slide_langgraph(prs)
    slide_subagent(prs)
    slide_persistence(prs)
    slide_future(prs)
    slide_closing(prs)

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"wrote {path}  ({path.stat().st_size / 1024:.1f} KiB)")
